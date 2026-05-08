"""Generate the Miinto Partner Vision SUMMIT presentation."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUTPUT = Path.home() / "projects/miinto-partner-onboarding/docs/Miinto-Partner-Vision-SUMMIT.pptx"

NAVY = RGBColor(0x1A, 0x36, 0x5D)
ACCENT = RGBColor(0xC9, 0xA2, 0x27)  # warm gold accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF5, 0xF7, 0xFA)
DARK_TEXT = RGBColor(0x1F, 0x2A, 0x3D)
MUTED = RGBColor(0x5A, 0x67, 0x85)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    set_fill(shape, color)
    return shape


def add_text(slide, left, top, width, height, text, *, size=18, bold=False,
             color=DARK_TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_header(slide, title, slide_num, total=14):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.0), NAVY)
    # accent stripe
    add_rect(slide, 0, Inches(1.0), SLIDE_W, Inches(0.06), ACCENT)
    add_text(slide, Inches(0.6), Inches(0.22), Inches(11), Inches(0.6),
             title, size=28, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(11.6), Inches(0.32), Inches(1.5), Inches(0.4),
             f"{slide_num} / {total}", size=12, color=WHITE,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_footer(slide, text="MIINTO  |  CONFIDENTIAL  |  Leadership Summit Q2 2026"):
    add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.35), LIGHT)
    add_text(slide, Inches(0.6), Inches(7.18), Inches(12.1), Inches(0.3),
             text, size=9, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)


def add_bullets(slide, items, *, left, top, width, height,
                size=20, line_spacing=1.35, color=DARK_TEXT, bold_first=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(8)
        # bullet marker
        marker = p.add_run()
        marker.text = "▸  "
        marker.font.name = "Calibri"
        marker.font.size = Pt(size)
        marker.font.bold = True
        marker.font.color.rgb = ACCENT
        # body
        run = p.add_run()
        run.text = item
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.bold = bold_first and idx == 0
        run.font.color.rgb = color
    return box


# ---------------------------------------------------------------------------
# Build deck
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


# Slide 1 — Title
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
add_rect(s, Inches(0.6), Inches(2.1), Inches(0.15), Inches(2.2), ACCENT)
add_text(s, Inches(1.0), Inches(2.0), Inches(11), Inches(1.4),
         "Miinto Partner Self-Service Vision",
         size=54, bold=True, color=WHITE)
add_text(s, Inches(1.0), Inches(3.4), Inches(11), Inches(0.7),
         "Leadership Summit — Q2 2026",
         size=26, color=ACCENT)
add_text(s, Inches(1.0), Inches(4.3), Inches(11), Inches(0.7),
         "From 600 to 6,000 partners. 18 months. One platform.",
         size=18, color=WHITE)
add_rect(s, Inches(1.0), Inches(6.6), Inches(2.0), Inches(0.04), ACCENT)
add_text(s, Inches(1.0), Inches(6.7), Inches(11), Inches(0.4),
         "CONFIDENTIAL", size=12, bold=True, color=WHITE)
add_text(s, Inches(11.0), Inches(6.7), Inches(1.8), Inches(0.4),
         "MIINTO", size=12, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)


# Slide 2 — The Imperative
s = prs.slides.add_slide(blank)
add_header(s, "The Imperative", 2)
add_bullets(s, [
    "Today: 600 partners, manual onboarding, 2-4 week wait",
    "Target: 6,000 partners by December 2027",
    "We cannot hire 10x more staff — we must go self-service",
], left=Inches(0.8), top=Inches(1.7), width=Inches(11.7), height=Inches(4.5),
   size=24)
# pull-quote panel
add_rect(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.2), LIGHT)
add_rect(s, Inches(0.8), Inches(5.6), Inches(0.12), Inches(1.2), ACCENT)
add_text(s, Inches(1.1), Inches(5.7), Inches(11.2), Inches(1.0),
         "10x growth in 18 months requires a 10x change in how we operate.",
         size=18, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s)


# Slide 3 — The Problem
s = prs.slides.add_slide(blank)
add_header(s, "The Problem", 3)
add_text(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.5),
         "Today's onboarding is slow, manual, and brittle:",
         size=18, color=MUTED)
add_bullets(s, [
    "Email back-and-forth to collect documents",
    "Print, sign, scan contracts",
    "Multiple calls with our team",
    "Days of waiting between each step",
], left=Inches(0.8), top=Inches(2.2), width=Inches(11.7), height=Inches(3.5),
   size=22)
# result callout
add_rect(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.3), NAVY)
add_text(s, Inches(1.1), Inches(5.5), Inches(11.2), Inches(1.3),
         "Result: Great boutiques give up.",
         size=26, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s)


# Slide 4 — Competitive Urgency
s = prs.slides.add_slide(blank)
add_header(s, "Competitive Urgency", 4)
# three competitor cards
labels = [
    ("ZALANDO", "Ahead of us today.\nMature partner platform."),
    ("FARFETCH", "Investing heavily\nin partner experience."),
    ("ABOUT YOU", "Doubling down on\nself-service tooling."),
]
card_w = Inches(3.9)
card_h = Inches(2.4)
gap = Inches(0.15)
total = card_w * 3 + gap * 2
start = (SLIDE_W - total) / 2
for i, (name, desc) in enumerate(labels):
    left = start + (card_w + gap) * i
    add_rect(s, left, Inches(1.7), card_w, card_h, LIGHT)
    add_rect(s, left, Inches(1.7), card_w, Inches(0.5), NAVY)
    add_text(s, left, Inches(1.7), card_w, Inches(0.5),
             name, size=16, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, left + Inches(0.2), Inches(2.35), card_w - Inches(0.4),
             Inches(1.6), desc, size=14, color=DARK_TEXT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_rect(s, Inches(0.8), Inches(4.7), Inches(11.7), Inches(2.0), NAVY)
add_rect(s, Inches(0.8), Inches(4.7), Inches(0.15), Inches(2.0), ACCENT)
add_text(s, Inches(1.2), Inches(4.85), Inches(11.0), Inches(0.55),
         "This window is closing.",
         size=28, bold=True, color=ACCENT)
add_text(s, Inches(1.2), Inches(5.45), Inches(11.0), Inches(1.2),
         "Zalando is ahead. Farfetch and About You are investing.\nWe must act now to secure premium boutique supply.",
         size=16, color=WHITE)
add_footer(s)


# Slide 5 — The Vision
s = prs.slides.add_slide(blank)
add_header(s, "The Vision", 5)
add_text(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.0),
         "18 months. 10x partner growth. One modern platform.",
         size=28, bold=True, color=NAVY)
# stat tiles
stats = [
    ("2-4 weeks → 48 hours", "Onboarding speed"),
    ("600 → 6,000", "Partners by 2027"),
    ("Bottleneck → Strategic", "Team focus"),
]
tile_w = Inches(3.9)
tile_h = Inches(2.6)
gap = Inches(0.15)
total = tile_w * 3 + gap * 2
start = (SLIDE_W - total) / 2
for i, (big, small) in enumerate(stats):
    left = start + (tile_w + gap) * i
    add_rect(s, left, Inches(3.2), tile_w, tile_h, LIGHT)
    add_rect(s, left, Inches(3.2), tile_w, Inches(0.08), ACCENT)
    add_text(s, left + Inches(0.2), Inches(3.6), tile_w - Inches(0.4),
             Inches(1.4), big, size=22, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, left + Inches(0.2), Inches(5.0), tile_w - Inches(0.4),
             Inches(0.6), small, size=14, color=MUTED,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.6),
         "Team capacity no longer the bottleneck — partner count is.",
         size=16, color=DARK_TEXT, align=PP_ALIGN.CENTER)
add_footer(s)


# Slide 6 — The Companion App
s = prs.slides.add_slide(blank)
add_header(s, "The Companion App", 6)
modules = [
    ("APPLY", "Submit, verify VAT, track status"),
    ("SIGN", "Digital contracts, EU-compliant"),
    ("LEARN", "Self-paced training"),
    ("SELL", "Orders, inventory, returns dashboard"),
    ("GET HELP", "Knowledge base, no phone calls"),
    ("CONFIGURE", "API keys, webhooks, integrations"),
]
cols = 3
rows = 2
m_w = Inches(4.0)
m_h = Inches(2.4)
gap_x = Inches(0.15)
gap_y = Inches(0.15)
grid_w = m_w * cols + gap_x * (cols - 1)
start_x = (SLIDE_W - grid_w) / 2
start_y = Inches(1.6)
for idx, (name, desc) in enumerate(modules):
    r = idx // cols
    c = idx % cols
    left = start_x + (m_w + gap_x) * c
    top = start_y + (m_h + gap_y) * r
    add_rect(s, left, top, m_w, m_h, LIGHT)
    add_rect(s, left, top, Inches(0.12), m_h, ACCENT)
    add_text(s, left + Inches(0.3), top + Inches(0.3), m_w - Inches(0.5),
             Inches(0.6), name, size=18, bold=True, color=NAVY)
    add_text(s, left + Inches(0.3), top + Inches(1.0), m_w - Inches(0.5),
             Inches(1.3), desc, size=14, color=DARK_TEXT)
add_footer(s)


# Slide 7 — The 48-Hour Journey
s = prs.slides.add_slide(blank)
add_header(s, "The 48-Hour Journey", 7)
steps = [
    ("HOUR 1", "Apply & VAT auto-verify"),
    ("HOUR 8", "Upload docs, KYC, digital sign"),
    ("DAY 2", "Connect Shopify, sandbox test"),
    ("DAY 3+", "Orders flow, go live"),
]
# horizontal timeline
tl_top = Inches(3.3)
tl_left = Inches(1.0)
tl_right = Inches(12.3)
add_rect(s, tl_left, tl_top + Inches(0.5), tl_right - tl_left, Inches(0.06), NAVY)
n = len(steps)
for i, (label, desc) in enumerate(steps):
    cx = tl_left + (tl_right - tl_left) * i / (n - 1)
    # node
    node = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              cx - Inches(0.25), tl_top + Inches(0.34),
                              Inches(0.5), Inches(0.5))
    set_fill(node, ACCENT)
    # label above
    add_text(s, cx - Inches(1.5), tl_top - Inches(0.9), Inches(3.0),
             Inches(0.5), label, size=18, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER)
    # desc below
    add_text(s, cx - Inches(1.5), tl_top + Inches(1.1), Inches(3.0),
             Inches(1.4), desc, size=14, color=DARK_TEXT,
             align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.6),
         "From application to first order in under 72 hours.",
         size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_footer(s)


# Slide 8 — Business Impact
s = prs.slides.add_slide(blank)
add_header(s, "Business Impact", 8)
rows_data = [
    ("Onboarding", "2-4 weeks", "48 hours"),
    ("Capacity", "~20 partners/month", "500+/month"),
    ("Total partners", "600", "6,000"),
    ("Team", "Bottleneck", "Strategic"),
]
tbl_left = Inches(0.8)
tbl_top = Inches(1.7)
col_w = [Inches(3.6), Inches(4.0), Inches(4.1)]
row_h = Inches(0.75)
# header row
headers = ["Dimension", "Today", "Target"]
x = tbl_left
for i, h in enumerate(headers):
    add_rect(s, x, tbl_top, col_w[i], row_h, NAVY)
    add_text(s, x + Inches(0.2), tbl_top, col_w[i] - Inches(0.3), row_h,
             h, size=16, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    x += col_w[i]
# body rows
for r, row in enumerate(rows_data):
    y = tbl_top + row_h * (r + 1)
    bg = LIGHT if r % 2 == 0 else WHITE
    x = tbl_left
    for i, val in enumerate(row):
        add_rect(s, x, y, col_w[i], row_h, bg)
        bold = i == 2
        color = NAVY if i == 2 else DARK_TEXT
        size = 16 if i == 0 else 16
        add_text(s, x + Inches(0.2), y, col_w[i] - Inches(0.3), row_h,
                 val, size=size, bold=bold or i == 0, color=color,
                 anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[i]

add_rect(s, Inches(0.8), Inches(5.4), Inches(11.7), Inches(1.4), LIGHT)
add_rect(s, Inches(0.8), Inches(5.4), Inches(0.12), Inches(1.4), ACCENT)
add_text(s, Inches(1.1), Inches(5.4), Inches(11.3), Inches(1.4),
         "Capacity scales 25x without scaling headcount 25x.",
         size=20, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s)


# Slide 9 — Financial Opportunity
s = prs.slides.add_slide(blank)
add_header(s, "Financial Opportunity", 9)
fin = [
    ("EUR 150k", "Average partner GMV / year"),
    ("5,400", "Incremental partners"),
    ("EUR 810M+", "Incremental annual GMV"),
]
tile_w = Inches(3.9)
tile_h = Inches(3.2)
gap = Inches(0.15)
total = tile_w * 3 + gap * 2
start = (SLIDE_W - total) / 2
for i, (big, small) in enumerate(fin):
    left = start + (tile_w + gap) * i
    add_rect(s, left, Inches(1.9), tile_w, tile_h, NAVY)
    add_rect(s, left, Inches(1.9), tile_w, Inches(0.1), ACCENT)
    add_text(s, left + Inches(0.2), Inches(2.4), tile_w - Inches(0.4),
             Inches(1.4), big, size=40, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, left + Inches(0.2), Inches(3.9), tile_w - Inches(0.4),
             Inches(1.0), small, size=14, color=ACCENT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.5),
         "Indicative figures — Finance to validate.",
         size=14, color=MUTED, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.6),
         "A platform investment with a category-defining return.",
         size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_footer(s)


# Slide 10 — Roadmap
s = prs.slides.add_slide(blank)
add_header(s, "Roadmap", 10)
phases = [
    ("Q2 2026", "FOUNDATION", "E-sign, KYC, wizard, sandbox"),
    ("Q3 2026", "SELF-SERVICE", "Training, integrations, support"),
    ("Q4 2026 – Q2 2027", "SCALE", "Mobile app, analytics"),
    ("Q3 – Q4 2027", "GROWTH", "EU expansion"),
]
ph_w = Inches(2.95)
ph_h = Inches(4.4)
gap = Inches(0.1)
total = ph_w * 4 + gap * 3
start = (SLIDE_W - total) / 2
for i, (when, name, desc) in enumerate(phases):
    left = start + (ph_w + gap) * i
    add_rect(s, left, Inches(1.7), ph_w, ph_h, LIGHT)
    add_rect(s, left, Inches(1.7), ph_w, Inches(0.6), NAVY)
    add_text(s, left, Inches(1.7), ph_w, Inches(0.6), when,
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, left + Inches(0.2), Inches(2.5), ph_w - Inches(0.4),
             Inches(0.7), name, size=20, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER)
    add_text(s, left + Inches(0.25), Inches(3.4), ph_w - Inches(0.5),
             Inches(2.5), desc, size=14, color=DARK_TEXT,
             align=PP_ALIGN.CENTER)
add_footer(s)


# Slide 11 — NTS Synergy
s = prs.slides.add_slide(blank)
add_header(s, "NTS Synergy", 11)
# equation visual
eq_top = Inches(2.2)
eq_h = Inches(2.0)
box_w = Inches(3.4)
gap = Inches(0.4)
plus_w = Inches(0.6)
eq_w = Inches(0.6)
result_w = Inches(3.6)
total_w = box_w + plus_w + box_w + eq_w + result_w + gap * 4
start = (SLIDE_W - total_w) / 2

x = start
add_rect(s, x, eq_top, box_w, eq_h, LIGHT)
add_rect(s, x, eq_top, box_w, Inches(0.08), ACCENT)
add_text(s, x, eq_top, box_w, eq_h, "Companion App",
         size=20, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
x += box_w + gap
add_text(s, x, eq_top, plus_w, eq_h, "+", size=44, bold=True,
         color=ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
x += plus_w + gap
add_rect(s, x, eq_top, box_w, eq_h, LIGHT)
add_rect(s, x, eq_top, box_w, Inches(0.08), ACCENT)
add_text(s, x, eq_top, box_w, eq_h, "NTS Migration",
         size=20, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
x += box_w + gap
add_text(s, x, eq_top, eq_w, eq_h, "=", size=44, bold=True,
         color=ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
x += eq_w + gap
add_rect(s, x, eq_top, result_w, eq_h, NAVY)
add_text(s, x, eq_top, result_w, eq_h, "ONE PROJECT",
         size=22, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_bullets(s, [
    "Early access for partners already on NTS",
    "By end of 2026: one modern system",
], left=Inches(2.0), top=Inches(5.0), width=Inches(9.3), height=Inches(2.0),
   size=20)
add_footer(s)


# Slide 12 — Key Risks
s = prs.slides.add_slide(blank)
add_header(s, "Key Risks", 12)
risks = [
    ("Integration complexity", "HIGH", "White-glove support",
     RGBColor(0xC0, 0x39, 0x2B)),
    ("Market variation", "MEDIUM", "Advisory board",
     RGBColor(0xD9, 0x82, 0x2B)),
    ("Adoption resistance", "MEDIUM", "Success manager",
     RGBColor(0xD9, 0x82, 0x2B)),
]
tbl_left = Inches(0.8)
tbl_top = Inches(1.7)
col_w = [Inches(4.5), Inches(2.0), Inches(5.2)]
row_h = Inches(0.75)
# header
headers = ["Risk", "Severity", "Mitigation"]
x = tbl_left
for i, h in enumerate(headers):
    add_rect(s, x, tbl_top, col_w[i], row_h, NAVY)
    add_text(s, x + Inches(0.2), tbl_top, col_w[i] - Inches(0.3), row_h,
             h, size=16, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    x += col_w[i]
# body
for r, (risk, sev, mit, sev_color) in enumerate(risks):
    y = tbl_top + row_h * (r + 1)
    bg = LIGHT if r % 2 == 0 else WHITE
    x = tbl_left
    add_rect(s, x, y, col_w[0], row_h, bg)
    add_text(s, x + Inches(0.2), y, col_w[0] - Inches(0.3), row_h,
             risk, size=16, bold=True, color=DARK_TEXT, anchor=MSO_ANCHOR.MIDDLE)
    x += col_w[0]
    add_rect(s, x, y, col_w[1], row_h, bg)
    # severity pill
    pill_w = Inches(1.4)
    pill_h = Inches(0.45)
    pill_x = x + (col_w[1] - pill_w) / 2
    pill_y = y + (row_h - pill_h) / 2
    pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              pill_x, pill_y, pill_w, pill_h)
    set_fill(pill, sev_color)
    add_text(s, pill_x, pill_y, pill_w, pill_h, sev,
             size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += col_w[1]
    add_rect(s, x, y, col_w[2], row_h, bg)
    add_text(s, x + Inches(0.2), y, col_w[2] - Inches(0.3), row_h,
             mit, size=15, color=DARK_TEXT, anchor=MSO_ANCHOR.MIDDLE)

add_rect(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.2), LIGHT)
add_rect(s, Inches(0.8), Inches(5.6), Inches(0.12), Inches(1.2), ACCENT)
add_text(s, Inches(1.1), Inches(5.6), Inches(11.3), Inches(1.2),
         "Each risk has a named owner and a defined mitigation path.",
         size=18, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s)


# Slide 13 — What We Need
s = prs.slides.add_slide(blank)
add_header(s, "What We Need", 13)
asks = [
    ("BUDGET", "EUR [X] over 18 months"),
    ("TEAM", "[Y] dedicated FTE"),
    ("KICKOFF", "Q2 2026, aligned with NTS"),
    ("SPONSOR", "Named executive sponsor"),
]
cols = 2
m_w = Inches(5.9)
m_h = Inches(2.4)
gap_x = Inches(0.2)
gap_y = Inches(0.2)
grid_w = m_w * cols + gap_x * (cols - 1)
start_x = (SLIDE_W - grid_w) / 2
start_y = Inches(1.6)
for idx, (name, desc) in enumerate(asks):
    r = idx // cols
    c = idx % cols
    left = start_x + (m_w + gap_x) * c
    top = start_y + (m_h + gap_y) * r
    add_rect(s, left, top, m_w, m_h, LIGHT)
    add_rect(s, left, top, Inches(0.15), m_h, ACCENT)
    add_text(s, left + Inches(0.4), top + Inches(0.3), m_w - Inches(0.6),
             Inches(0.7), name, size=18, bold=True, color=ACCENT)
    add_text(s, left + Inches(0.4), top + Inches(1.0), m_w - Inches(0.6),
             Inches(1.3), desc, size=22, bold=True, color=NAVY)
add_footer(s)


# Slide 14 — The Ask
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
add_rect(s, Inches(0.6), Inches(0.9), Inches(0.15), Inches(1.2), ACCENT)
add_text(s, Inches(1.0), Inches(0.8), Inches(11), Inches(1.4),
         "The Ask", size=54, bold=True, color=WHITE)
add_text(s, Inches(1.0), Inches(1.8), Inches(11), Inches(0.6),
         "Why this, why now",
         size=20, color=ACCENT)

# four big statements
points = [
    "Grow from 600 to 6,000 partners",
    "Growth limited only by boutique count",
    "Zalando ahead — window closes soon",
]
top = Inches(2.9)
for i, txt in enumerate(points):
    y = top + Inches(0.7) * i
    add_text(s, Inches(1.1), y, Inches(0.5), Inches(0.6),
             "▸", size=24, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.6), y, Inches(11), Inches(0.6),
             txt, size=22, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

add_rect(s, Inches(1.0), Inches(5.6), Inches(11.3), Inches(1.2), ACCENT)
add_text(s, Inches(1.0), Inches(5.6), Inches(11.3), Inches(1.2),
         "Let's build this together.",
         size=32, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.4),
         "MIINTO  |  CONFIDENTIAL  |  Leadership Summit Q2 2026",
         size=10, color=WHITE, align=PP_ALIGN.CENTER)


OUTPUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
